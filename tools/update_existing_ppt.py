import pptx
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# 打开现有PPT
prs = pptx.Presentation('g:/iMato/scripts/MatuX_创业大赛商业计划书_完整版.pptx')

# 定义颜色
COLOR_PRIMARY = RGBColor(26, 35, 126)    # 深蓝色
COLOR_SECONDARY = RGBColor(0, 168, 150)   # 青色
COLOR_ACCENT = RGBColor(255, 107, 107)    # 红色
COLOR_WHITE = RGBColor(255, 255, 255)
COLOR_GRAY = RGBColor(102, 102, 102)
COLOR_BG_DARK = RGBColor(21, 32, 43)      # 午夜蓝
COLOR_BG_LIGHT = RGBColor(242, 244, 248) # 浅灰色

def add_title_slide(title, subtitle=""):
    """添加标题页"""
    slide_layout = prs.slide_layouts[5]  # 使用可用的布局
    slide = prs.slides.add_slide(slide_layout)
    
    # 添加标题
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(9)
    height = Inches(1)
    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = COLOR_PRIMARY
    title_para.alignment = PP_ALIGN.CENTER
    
    if subtitle:
        subtitle_box = slide.shapes.add_textbox(left, top + Inches(1.2), width, height)
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = subtitle
        subtitle_para = subtitle_frame.paragraphs[0]
        subtitle_para.font.size = Pt(24)
        subtitle_para.font.color.rgb = COLOR_SECONDARY
        subtitle_para.alignment = PP_ALIGN.CENTER
    
    return slide

def add_card_textbox(slide, left, top, width, height, title, content, bg_color, title_color=COLOR_WHITE):
    """添加卡片式文本框"""
    shape = slide.shapes.add_shape(1, left, top, width, height)  # 矩形
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.fill.background()
    
    text_frame = shape.text_frame
    text_frame.word_wrap = True
    text_frame.margin_left = Inches(0.15)
    text_frame.margin_right = Inches(0.15)
    text_frame.margin_top = Inches(0.1)
    text_frame.margin_bottom = Inches(0.1)
    
    # 添加标题
    p = text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = title_color
    p.alignment = PP_ALIGN.LEFT
    
    # 添加内容
    if content:
        p = text_frame.add_paragraph()
        p.text = content
        p.font.size = Pt(11)
        p.font.color.rgb = COLOR_WHITE
        p.alignment = PP_ALIGN.LEFT
        p.space_before = Pt(6)
    
    return shape

def add_two_column_slide(title, left_title, left_content, right_title, right_content):
    """添加双列幻灯片"""
    slide = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide)
    
    # 添加主标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = COLOR_PRIMARY
    title_para.alignment = PP_ALIGN.CENTER
    
    # 左列
    add_card_textbox(slide, Inches(0.5), Inches(1.6), Inches(4.2), Inches(4.5), 
                    left_title, left_content, COLOR_PRIMARY)
    
    # 右列
    add_card_textbox(slide, Inches(5.3), Inches(1.6), Inches(4.2), Inches(4.5), 
                    right_title, right_content, COLOR_SECONDARY)
    
    return slide

# ==================== 第1页：K12开源管理系统 ====================
slide = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide)

# 标题
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
title_frame = title_box.text_frame
title_frame.text = "K12开源管理系统"
title_para = title_frame.paragraphs[0]
title_para.font.size = Pt(36)
title_para.font.bold = True
title_para.font.color.rgb = COLOR_PRIMARY
title_para.alignment = PP_ALIGN.CENTER

# 卡片1：产品定位
add_card_textbox(slide, Inches(0.5), Inches(1.6), Inches(4.2), Inches(1.8),
                "📚 产品定位",
                "为教育机构、校园AI教学、机器人培训机构（面向K12）提供开源免费管理系统，带AI教师、基础入门课件，打通学生、家长、机构、学校，建立统一的课程体系。",
                COLOR_PRIMARY)

# 卡片2：目标用户
add_card_textbox(slide, Inches(5.3), Inches(1.6), Inches(4.2), Inches(1.8),
                "🎯 目标用户",
                "教育机构 | 学校 | 教师 | 学生 | 家长 | 教育局",
                COLOR_SECONDARY)

# 卡片3：核心功能
add_card_textbox(slide, Inches(0.5), Inches(3.6), Inches(4.2), Inches(2.5),
                "⭐ 核心功能",
                "• 统一课程体系（robotics/programming/ai_fundamentals/project_based）\n• AI教师系统（智能问答、代码批改、学习推荐）\n• 四角色Dashboard（教师、机构、学校、教育局）\n• 多端协同（Web、移动、管理端）",
                COLOR_BG_DARK, COLOR_WHITE)

# 卡片4：开源策略
add_card_textbox(slide, Inches(5.3), Inches(3.6), Inches(4.2), Inches(2.5),
                "🔓 开源策略",
                "• 协议：AGPL-3.0\n• 包含：完整代码、基础课程、API文档\n• 商业化：免费版 + 专业版 + 企业版",
                COLOR_BG_LIGHT, COLOR_PRIMARY)

# ==================== 第2页：AI教师系统 ====================
slide = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide)

# 标题
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
title_frame = title_box.text_frame
title_frame.text = "AI教师系统"
title_para = title_frame.paragraphs[0]
title_para.font.size = Pt(36)
title_para.font.bold = True
title_para.font.color.rgb = COLOR_PRIMARY
title_para.alignment = PP_ALIGN.CENTER

# 卡片
cards = [
    ("💬 智能问答", "基于课程内容回答学生问题\n7×24小时在线答疑", 0, 1.6),
    ("✅ 代码批改", "智能分析学生代码并给出反馈\n自动评分和改进建议", 4.8, 1.6),
    ("🎯 项目指导", "分步骤指导学生完成实践项目\n实时问题解答", 0, 3.6),
    ("📊 学习推荐", "个性化学习路径推荐\n基于学习历史和兴趣", 4.8, 3.6),
]

for title, content, x, y in cards:
    add_card_textbox(slide, Inches(x), Inches(y), Inches(4.2), Inches(1.5), 
                    title, content, COLOR_PRIMARY if x == 0 else COLOR_SECONDARY)

# Token消耗说明
add_card_textbox(slide, Inches(0.5), Inches(5.3), Inches(8.5), Inches(1.3),
                "💰 Token消耗参考",
                "智能问答: 300-1500 tokens | 代码批改: 800-3500 tokens | 学习推荐: 700-2500 tokens | 项目指导: 800-3000 tokens",
                COLOR_BG_LIGHT, COLOR_PRIMARY)

# ==================== 第3页：职校创客教育（第1部分）====================
slide = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide)

# 标题
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
title_frame = title_box.text_frame
title_frame.text = "职校创客教育系统（1/3）"
title_para = title_frame.paragraphs[0]
title_para.font.size = Pt(36)
title_para.font.bold = True
title_para.font.color.rgb = COLOR_PRIMARY
title_para.alignment = PP_ALIGN.CENTER

# 核心价值
add_card_textbox(slide, Inches(0.5), Inches(1.6), Inches(8.5), Inches(0.8),
                "💡 核心价值",
                "连接学校教育与企业实际需求，激发学生创新能力和实践能力，支持创客赛事和就业对接",
                COLOR_BG_LIGHT, COLOR_PRIMARY)

# 创客项目孵化
add_card_textbox(slide, Inches(0.5), Inches(2.6), Inches(4.2), Inches(2.2),
                "🚀 创客项目孵化",
                "• 机器人项目（智能小车、机械臂、无人机）\n• 物联网项目（智能家居、环境监测）\n• 软件项目（移动应用、Web应用）\n• 混合项目（软硬件结合）",
                COLOR_PRIMARY)

# 企业需求对接
add_card_textbox(slide, Inches(5.3), Inches(2.6), Inches(4.2), Inches(2.2),
                "🏢 企业需求对接",
                "• 技术开发需求（小程序、网站建设）\n• 产品设计需求（UI/UX、原型设计）\n• 创新研究需求（市场调研、技术预研）\n• 校企合作项目（联合研发、实习实训）",
                COLOR_SECONDARY)

# 流程说明
add_card_textbox(slide, Inches(0.5), Inches(5.0), Inches(4.2), Inches(1.6),
                "📋 项目流程",
                "创意提交 → 导师审核 → 项目立项 → 资源分配 → 开发实施 → 阶段评审 → 成果展示",
                COLOR_BG_DARK, COLOR_WHITE)

add_card_textbox(slide, Inches(5.3), Inches(5.0), Inches(4.2), Inches(1.6),
                "🤝 对接流程",
                "企业发布需求 → 学校审核 → 学生组队竞标 → 企业选择 → 签订协议 → 项目实施 → 验收交付",
                COLOR_BG_LIGHT, COLOR_PRIMARY)

# ==================== 第4页：职校创客教育（第2部分）====================
slide = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide)

# 标题
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
title_frame = title_box.text_frame
title_frame.text = "职校创客教育系统（2/3）"
title_para = title_frame.paragraphs[0]
title_para.font.size = Pt(36)
title_para.font.bold = True
title_para.font.color.rgb = COLOR_PRIMARY
title_para.alignment = PP_ALIGN.CENTER

# 市场需求推荐
add_card_textbox(slide, Inches(0.5), Inches(1.6), Inches(4.2), Inches(2.2),
                "🤖 市场需求推荐系统",
                "基于AI算法的智能匹配：\n• 学生技能匹配度（技能标签、技术栈）\n• 历史项目相似度\n• 学习兴趣偏好\n• 同类学生的参与记录",
                COLOR_PRIMARY)

# 创客赛事管理
add_card_textbox(slide, Inches(5.3), Inches(1.6), Inches(4.2), Inches(2.2),
                "🏆 创客赛事管理",
                "赛事分类：\n• 校内赛事（学期项目展示、创新大赛）\n• 区域赛事（市级、省级竞赛）\n• 全国赛事（中国创客大赛、机器人锦标赛）\n• 国际赛事（FIRST、RoboCup）",
                COLOR_SECONDARY)

# 成果展示
add_card_textbox(slide, Inches(0.5), Inches(4.0), Inches(4.2), Inches(2.6),
                "💼 成果展示与就业对接",
                "• 个人成果库：项目作品、竞赛成就、技能认证\n• 竞赛成果：获奖记录、证书、媒体报道\n• 能力矩阵：技能雷达图、成长轨迹\n• 就业推荐：基于项目经验的企业职位推荐",
                COLOR_BG_DARK, COLOR_WHITE)

# 就业对接
add_card_textbox(slide, Inches(5.3), Inches(4.0), Inches(4.2), Inches(2.6),
                "🎓 就业对接功能",
                "• 智能职位推荐\n• 技能匹配度分析\n• 项目作品集展示\n• 企业直聘通道\n• 实习机会对接",
                COLOR_BG_LIGHT, COLOR_PRIMARY)

# ==================== 第5页：职校创客教育（第3部分）====================
slide = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide)

# 标题
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
title_frame = title_box.text_frame
title_frame.text = "职校创客教育系统（3/3）"
title_para = title_frame.paragraphs[0]
title_para.font.size = Pt(36)
title_para.font.bold = True
title_para.font.color.rgb = COLOR_PRIMARY
title_para.alignment = PP_ALIGN.CENTER

# 实施计划
phases = [
    ("第一阶段", "基础功能开发", "2-3个月", "创客项目管理、团队协作、项目展示"),
    ("第二阶段", "企业需求对接", "1-2个月", "需求发布、投标竞标、合同管理"),
    ("第三阶段", "赛事管理", "1-2个月", "赛事发布、团队报名、成果评审"),
    ("第四阶段", "智能推荐", "1-2个月", "AI需求推荐、成果库、就业对接"),
    ("第五阶段", "运营优化", "持续", "数据分析、用户反馈、功能迭代"),
]

y_pos = 1.6
for idx, (phase, name, duration, tasks) in enumerate(phases):
    bg_color = COLOR_PRIMARY if idx % 2 == 0 else COLOR_SECONDARY
    add_card_textbox(slide, Inches(0.5), Inches(y_pos), Inches(1.5), Inches(0.9), 
                    phase, "", bg_color)
    
    add_card_textbox(slide, Inches(2.2), Inches(y_pos), Inches(2.5), Inches(0.9), 
                    name, duration, COLOR_BG_LIGHT, COLOR_PRIMARY)
    
    add_card_textbox(slide, Inches(4.9), Inches(y_pos), Inches(4.1), Inches(0.9), 
                    "任务", tasks, COLOR_BG_DARK, COLOR_WHITE)
    
    y_pos += 1.0

# 总体时间
add_card_textbox(slide, Inches(0.5), Inches(6.7), Inches(8.5), Inches(0.9),
                "⏱️ 总体时间规划",
                "预计总工期：7-11个月 | 里程碑：第3个月（企业需求上线）、第5个月（赛事功能上线）、第8个月（AI推荐上线）",
                COLOR_BG_LIGHT, COLOR_PRIMARY)

# ==================== 第6页：收费模型设计 ====================
slide = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide)

# 标题
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
title_frame = title_box.text_frame
title_frame.text = "收费模型：免费课程体系 + AI课程（订阅费 + Token消耗）"
title_para = title_frame.paragraphs[0]
title_para.font.size = Pt(32)
title_para.font.bold = True
title_para.font.color.rgb = COLOR_PRIMARY
title_para.alignment = PP_ALIGN.CENTER

# 免费版
add_card_textbox(slide, Inches(0.5), Inches(1.6), Inches(2.8), Inches(3.0),
                "🆓 免费版",
                "• 价格：¥0\n• 适用：<50学生\n• 开源代码完整\n• 基础课程（Level 1-2）\n• 四角色基础Dashboard\n• 创客项目管理\n• ❌ 无AI功能",
                COLOR_BG_LIGHT, COLOR_PRIMARY)

# 专业版
add_card_textbox(slide, Inches(3.6), Inches(1.6), Inches(2.8), Inches(3.0),
                "⭐ 专业版",
                "• 价格：¥99/月 或 ¥999/年\n• 适用：50-500学生\n• ✅ AI教师（10K tokens/月）\n• 智能课程推荐\n• 学习路径AI规划\n• 代码自动批改（500次/月）\n• 企业需求推荐\n• 赛事智能匹配",
                COLOR_SECONDARY)

# 企业版
add_card_textbox(slide, Inches(6.7), Inches(1.6), Inches(2.8), Inches(3.0),
                "🏆 企业版",
                "• 价格：¥2999/月 或 ¥29999/年\n• 适用：>500学生\n• ✅ AI教师（100K tokens/月）\n• 专属技术支持（7×24）\n• 定制化功能开发\n• 数据导入/导出服务\n• 专属服务器部署\n• SLA保障（99.9%）",
                COLOR_PRIMARY)

# 核心优势
add_card_textbox(slide, Inches(0.5), Inches(4.8), Inches(8.5), Inches(1.0),
                "💎 核心优势",
                "• 开源免费降低使用门槛 • AI功能按需付费 • Token消耗透明可控 • 灵活订阅满足不同规模",
                COLOR_BG_DARK, COLOR_WHITE)

# 适用场景
add_card_textbox(slide, Inches(0.5), Inches(5.9), Inches(8.5), Inches(1.7),
                "🎯 适用场景",
                "个人教师/小型机构 → 免费版（满足基础需求）\n中型机构/学校 → 专业版（AI赋能教学）\n大型机构/教育局 → 企业版（完整解决方案）",
                COLOR_BG_LIGHT, COLOR_PRIMARY)

# ==================== 第7页：Token计费规则 ====================
slide = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide)

# 标题
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
title_frame = title_box.text_frame
title_frame.text = "Token计费规则与提醒策略"
title_para = title_frame.paragraphs[0]
title_para.font.size = Pt(36)
title_para.font.bold = True
title_para.font.color.rgb = COLOR_PRIMARY
title_para.alignment = PP_ALIGN.CENTER

# Token消耗表格
table_data = [
    ["功能", "输入Tokens", "输出Tokens", "单次消耗"],
    ["智能问答", "100-500", "200-1000", "300-1500"],
    ["代码批改", "500-2000", "300-1500", "800-3500"],
    ["学习推荐", "200-500", "500-2000", "700-2500"],
    ["项目指导", "300-1000", "500-2000", "800-3000"],
    ["课程生成", "500-1000", "1000-5000", "1500-6000"],
]

table = slide.shapes.add_table(len(table_data), 4, Inches(0.5), Inches(1.5), Inches(8.5), Inches(1.8))
table = table.table

# 设置表头
for col, text in enumerate(table_data[0]):
    cell = table.rows[0].cells[col]
    cell.text = text
    cell.fill.solid()
    cell.fill.fore_color.rgb = COLOR_PRIMARY
    paragraph = cell.text_frame.paragraphs[0]
    paragraph.font.bold = True
    paragraph.font.color.rgb = COLOR_WHITE
    paragraph.font.size = Pt(11)
    paragraph.alignment = PP_ALIGN.CENTER

# 设置表格内容
for row_idx, row_data in enumerate(table_data[1:], start=1):
    for col_idx, text in enumerate(row_data):
        cell = table.rows[row_idx].cells[col_idx]
        cell.text = text
        paragraph = cell.text_frame.paragraphs[0]
        paragraph.font.size = Pt(10)
        paragraph.alignment = PP_ALIGN.CENTER
        if row_idx % 2 == 1:
            cell.fill.solid()
            cell.fill.fore_color.rgb = COLOR_BG_LIGHT

# 计费规则
add_card_textbox(slide, Inches(0.5), Inches(3.5), Inches(4.2), Inches(1.8),
                "💰 计费规则",
                "• 免费版：¥0.1/1000 tokens\n• 专业版：含10K tokens，超出¥0.08/1000\n• 企业版：含100K tokens，超出¥0.06/1000",
                COLOR_SECONDARY)

# 提醒策略
add_card_textbox(slide, Inches(5.3), Inches(3.5), Inches(3.7), Inches(1.8),
                "🔔 提醒策略",
                "• 80%配额预警\n• 100%配额耗尽（暂停AI）\n• 到期前7/3/1天提醒\n• 未续费自动降级",
                COLOR_PRIMARY)

# 提醒策略详情
add_card_textbox(slide, Inches(0.5), Inches(5.5), Inches(8.5), Inches(2.1),
                "📊 配额管理与降级说明",
                "配额预警：系统自动发送邮件和内站通知，提示即将耗尽\n配额耗尽：暂停AI教师功能，建议购买额外配额或升级套餐\n到期提醒：多渠道提醒用户及时续费，避免服务中断\n功能降级：订阅过期后保留所有数据，但AI功能降级至免费版，数据可恢复",
                COLOR_BG_DARK, COLOR_WHITE)

# 保存文件
prs.save('g:/iMato/scripts/MatuX_创业大赛商业计划书_完整版.pptx')

print("✅ PPT已成功更新！")
print(f"共 {len(prs.slides)} 页")
print("新增内容：")
print("  - K12开源管理系统（1页）")
print("  - AI教师系统（1页）")
print("  - 职校创客教育系统（3页）")
print("  - 收费模型设计（1页）")
print("  - Token计费规则（1页）")