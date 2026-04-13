"""
文档处理服务
支持PDF、Markdown、PPTX等文档格式的处理和转换
"""

import io
import logging
import os
from typing import Any, Dict, List, Optional, Tuple

from PIL import Image
from docx import Document
import fitz  # PyMuPDF
import markdown
from pptx import Presentation
from sqlalchemy.orm import Session

from models.multimedia import DocumentFormat, MediaType, MultimediaResource
from services.multimedia_service import StorageConfig

logger = logging.getLogger(__name__)


class DocumentProcessingService:
    """文档处理服务类"""

    def __init__(self, db: Session):
        self.db = db
        self.storage_config = StorageConfig()
        self.max_image_size = (800, 600)  # 缩略图最大尺寸

    def process_document(
        self,
        resource_id: int,
        convert_to_pdf: bool = True,
        generate_thumbnails: bool = True,
        extract_text: bool = True,
    ) -> Dict[str, Any]:
        """
        处理文档资源

        Args:
            resource_id: 资源ID
            convert_to_pdf: 是否转换为PDF
            generate_thumbnails: 是否生成缩略图
            extract_text: 是否提取文本

        Returns:
            Dict: 处理结果
        """
        try:
            # 获取资源信息
            resource = (
                self.db.query(MultimediaResource)
                .filter(MultimediaResource.id == resource_id)
                .first()
            )

            if not resource:
                raise ValueError(f"资源不存在: {resource_id}")

            if resource.media_type != MediaType.DOCUMENT:
                raise ValueError("资源不是文档类型")

            # 下载文档文件
            document_data = self._download_document_file(resource)

            # 检测文档格式
            detected_format = self._detect_document_format(
                resource.file_name, document_data
            )
            resource.document_format = detected_format

            # 处理结果
            processing_results = {
                "text_extracted": False,
                "pdf_generated": False,
                "thumbnails_generated": False,
                "page_count": 0,
                "word_count": 0,
            }

            # 提取文本
            if extract_text:
                text_content, word_count = self._extract_text(
                    document_data, detected_format
                )
                processing_results["text_extracted"] = True
                processing_results["word_count"] = word_count

                # 保存文本内容到custom_metadata
                if not resource.custom_metadata:
                    resource.custom_metadata = {}
                resource.custom_metadata["extracted_text"] = text_content[
                    :10000
                ]  # 限制长度

            # 转换为PDF
            if convert_to_pdf and detected_format != DocumentFormat.PDF:
                pdf_data = self._convert_to_pdf(document_data, detected_format)
                if pdf_data:
                    pdf_url = self._save_processed_file(pdf_data, resource, "pdf")
                    resource.processed_url = pdf_url
                    processing_results["pdf_generated"] = True

            # 生成缩略图
            if generate_thumbnails:
                thumbnail_urls = self._generate_thumbnails(
                    document_data, detected_format, resource
                )
                processing_results["thumbnails_generated"] = len(thumbnail_urls) > 0
                processing_results["thumbnail_urls"] = thumbnail_urls

            # 更新资源信息
            if processing_results["text_extracted"]:
                resource.page_count = processing_results.get("page_count", 1)
                resource.word_count = processing_results["word_count"]

            self.db.commit()

            logger.info(f"文档处理完成: {resource_id}")

            return {
                "success": True,
                "resource_id": resource_id,
                "processing_results": processing_results,
            }

        except Exception as e:
            logger.error(f"文档处理失败: {e}")
            raise

    def _download_document_file(self, resource: MultimediaResource) -> bytes:
        """下载文档文件"""
        if self.storage_config.storage_type == "s3":
            return self._download_from_s3(resource.original_url)
        else:
            return self._read_local_file(resource.original_url)

    def _download_from_s3(self, file_url: str) -> bytes:
        """从S3下载文件"""

        if file_url.startswith("https://"):
            parts = file_url.split("/")
            bucket = parts[2].split(".")[0]
            key = "/".join(parts[3:])

            s3_client = self.storage_config.get_s3_client()
            response = s3_client.get_object(Bucket=bucket, Key=key)
            return response["Body"].read()
        else:
            raise ValueError("无效的S3 URL格式")

    def _read_local_file(self, file_url: str) -> bytes:
        """读取本地文件"""
        if file_url.startswith("/"):
            file_path = file_url[1:]
        else:
            file_path = file_url

        full_path = os.path.join(os.getcwd(), file_path)
        with open(full_path, "rb") as f:
            return f.read()

    def _detect_document_format(
        self, filename: str, file_data: bytes
    ) -> DocumentFormat:
        """检测文档格式"""
        extension = filename.lower().split(".")[-1] if "." in filename else ""

        # 根据扩展名判断
        format_mapping = {
            "pdf": DocumentFormat.PDF,
            "md": DocumentFormat.MARKDOWN,
            "markdown": DocumentFormat.MARKDOWN,
            "pptx": DocumentFormat.PPTX,
            "docx": DocumentFormat.DOCX,
            "txt": DocumentFormat.TXT,
            "html": DocumentFormat.HTML,
        }

        if extension in format_mapping:
            return format_mapping[extension]

        # 根据文件内容判断
        try:
            if file_data.startswith(b"%PDF"):
                return DocumentFormat.PDF
            elif b"<html" in file_data.lower() or b"<!doctype" in file_data.lower():
                return DocumentFormat.HTML
            elif file_data.startswith(b"PK"):  # ZIP格式，可能是DOCX或PPTX
                # 进一步检查内部结构
                import zipfile

                with zipfile.ZipFile(io.BytesIO(file_data)) as zf:
                    if "[Content_Types].xml" in zf.namelist():
                        if any(name.startswith("ppt/") for name in zf.namelist()):
                            return DocumentFormat.PPTX
                        elif any(name.startswith("word/") for name in zf.namelist()):
                            return DocumentFormat.DOCX
        except Exception:
            pass

        return DocumentFormat.TXT

    def _extract_text(
        self, document_data: bytes, doc_format: DocumentFormat
    ) -> Tuple[str, int]:
        """提取文档文本内容"""
        try:
            if doc_format == DocumentFormat.PDF:
                return self._extract_pdf_text(document_data)
            elif doc_format == DocumentFormat.DOCX:
                return self._extract_docx_text(document_data)
            elif doc_format == DocumentFormat.PPTX:
                return self._extract_pptx_text(document_data)
            elif doc_format == DocumentFormat.MARKDOWN:
                return self._extract_markdown_text(document_data)
            elif doc_format == DocumentFormat.HTML:
                return self._extract_html_text(document_data)
            elif doc_format == DocumentFormat.TXT:
                return self._extract_txt_text(document_data)
            else:
                return "", 0

        except Exception as e:
            logger.error(f"文本提取失败: {e}")
            return "", 0

    def _extract_pdf_text(self, pdf_data: bytes) -> Tuple[str, int]:
        """提取PDF文本"""
        try:
            pdf_document = fitz.open(stream=pdf_data, filetype="pdf")
            text_content = ""
            page_count = len(pdf_document)

            for page_num in range(page_count):
                page = pdf_document[page_num]
                text_content += page.get_text()

            word_count = len(text_content.split())
            pdf_document.close()

            return text_content, word_count

        except Exception as e:
            logger.error(f"PDF文本提取失败: {e}")
            raise

    def _extract_docx_text(self, docx_data: bytes) -> Tuple[str, int]:
        """提取DOCX文本"""
        try:
            doc = Document(io.BytesIO(docx_data))
            text_content = ""

            for paragraph in doc.paragraphs:
                text_content += paragraph.text + "\n"

            # 提取表格文本
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        text_content += cell.text + " "
                    text_content += "\n"

            word_count = len(text_content.split())
            return text_content, word_count

        except Exception as e:
            logger.error(f"DOCX文本提取失败: {e}")
            raise

    def _extract_pptx_text(self, pptx_data: bytes) -> Tuple[str, int]:
        """提取PPTX文本"""
        try:
            prs = Presentation(io.BytesIO(pptx_data))
            text_content = ""
            len(prs.slides)

            for slide_num, slide in enumerate(prs.slides):
                text_content += f"幻灯片 {slide_num + 1}:\n"

                # 提取幻灯片文本
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_content += shape.text + "\n"

                text_content += "\n"

            word_count = len(text_content.split())
            return text_content, word_count

        except Exception as e:
            logger.error(f"PPTX文本提取失败: {e}")
            raise

    def _extract_markdown_text(self, md_data: bytes) -> Tuple[str, int]:
        """提取Markdown文本"""
        try:
            md_text = md_data.decode("utf-8", errors="ignore")
            # 简单移除Markdown标记
            import re

            # 移除标题标记
            clean_text = re.sub(r"^#{1,6}\s+", "", md_text, flags=re.MULTILINE)
            # 移除强调标记
            clean_text = re.sub(r"[*_]{1,2}(.*?)[*_]{1,2}", r"\1", clean_text)
            # 移除代码块标记
            clean_text = re.sub(r"`([^`]+)`", r"\1", clean_text)
            # 移除链接
            clean_text = re.sub(r"\[([^\]]+)\]\([^\)]+\)", r"\1", clean_text)

            word_count = len(clean_text.split())
            return clean_text, word_count

        except Exception as e:
            logger.error(f"Markdown文本提取失败: {e}")
            raise

    def _extract_html_text(self, html_data: bytes) -> Tuple[str, int]:
        """提取HTML文本"""
        try:
            from bs4 import BeautifulSoup

            html_text = html_data.decode("utf-8", errors="ignore")
            soup = BeautifulSoup(html_text, "html.parser")
            clean_text = soup.get_text()
            word_count = len(clean_text.split())
            return clean_text, word_count

        except Exception as e:
            logger.error(f"HTML文本提取失败: {e}")
            raise

    def _extract_txt_text(self, txt_data: bytes) -> Tuple[str, int]:
        """提取TXT文本"""
        try:
            text_content = txt_data.decode("utf-8", errors="ignore")
            word_count = len(text_content.split())
            return text_content, word_count

        except Exception as e:
            logger.error(f"TXT文本提取失败: {e}")
            raise

    def _convert_to_pdf(
        self, document_data: bytes, doc_format: DocumentFormat
    ) -> Optional[bytes]:
        """将文档转换为PDF"""
        try:
            if doc_format == DocumentFormat.MARKDOWN:
                return self._markdown_to_pdf(document_data)
            elif doc_format == DocumentFormat.HTML:
                return self._html_to_pdf(document_data)
            elif doc_format == DocumentFormat.TXT:
                return self._txt_to_pdf(document_data)
            else:
                # 对于其他格式，尝试使用LibreOffice转换（需要额外配置）
                logger.warning(f"暂不支持 {doc_format} 格式转换为PDF")
                return None

        except Exception as e:
            logger.error(f"文档转换PDF失败: {e}")
            return None

    def _markdown_to_pdf(self, md_data: bytes) -> bytes:
        """Markdown转PDF"""
        try:
            # 先转换为HTML
            md_text = md_data.decode("utf-8", errors="ignore")
            html_content = markdown.markdown(md_text)

            # 使用WeasyPrint转换为PDF
            from weasyprint import HTML

            pdf_bytes = HTML(string=html_content).write_pdf()
            return pdf_bytes

        except Exception as e:
            logger.error(f"Markdown转PDF失败: {e}")
            raise

    def _html_to_pdf(self, html_data: bytes) -> bytes:
        """HTML转PDF"""
        try:
            html_text = html_data.decode("utf-8", errors="ignore")
            from weasyprint import HTML

            pdf_bytes = HTML(string=html_text).write_pdf()
            return pdf_bytes

        except Exception as e:
            logger.error(f"HTML转PDF失败: {e}")
            raise

    def _txt_to_pdf(self, txt_data: bytes) -> bytes:
        """TXT转PDF"""
        try:
            txt_text = txt_data.decode("utf-8", errors="ignore")

            # 简单的文本转PDF
            import io

            from reportlab.lib.pagesizes import letter
            from reportlab.pdfgen import canvas

            buffer = io.BytesIO()
            c = canvas.Canvas(buffer, pagesize=letter)
            width, height = letter

            # 设置字体和大小
            c.setFont("Helvetica", 12)

            # 分割文本为行
            lines = txt_text.split("\n")
            y_position = height - 50

            for line in lines:
                if y_position < 50:  # 需要新页面
                    c.showPage()
                    c.setFont("Helvetica", 12)
                    y_position = height - 50

                c.drawString(50, y_position, line[:100])  # 限制每行长度
                y_position -= 15

            c.save()
            pdf_bytes = buffer.getvalue()
            buffer.close()

            return pdf_bytes

        except Exception as e:
            logger.error(f"TXT转PDF失败: {e}")
            raise

    def _generate_thumbnails(
        self,
        document_data: bytes,
        doc_format: DocumentFormat,
        resource: MultimediaResource,
    ) -> List[str]:
        """生成文档缩略图"""
        try:
            thumbnail_urls = []

            if doc_format == DocumentFormat.PDF:
                thumbnail_urls = self._generate_pdf_thumbnails(document_data, resource)
            elif doc_format == DocumentFormat.PPTX:
                thumbnail_urls = self._generate_pptx_thumbnails(document_data, resource)
            elif doc_format in [
                DocumentFormat.DOCX,
                DocumentFormat.MARKDOWN,
                DocumentFormat.HTML,
                DocumentFormat.TXT,
            ]:
                # 为这些格式生成单个封面缩略图
                thumbnail_url = self._generate_cover_thumbnail(
                    document_data, doc_format, resource
                )
                if thumbnail_url:
                    thumbnail_urls.append(thumbnail_url)

            return thumbnail_urls

        except Exception as e:
            logger.error(f"缩略图生成失败: {e}")
            return []

    def _generate_pdf_thumbnails(
        self, pdf_data: bytes, resource: MultimediaResource
    ) -> List[str]:
        """生成PDF缩略图"""
        try:
            pdf_document = fitz.open(stream=pdf_data, filetype="pdf")
            thumbnail_urls = []

            # 为前几页生成缩略图
            pages_to_thumbnail = min(3, len(pdf_document))

            for page_num in range(pages_to_thumbnail):
                page = pdf_document[page_num]

                # 生成图片
                mat = fitz.Matrix(2.0, 2.0)  # 2倍分辨率
                pix = page.get_pixmap(matrix=mat)

                # 转换为PIL Image
                img_data = pix.tobytes("ppm")
                img = Image.open(io.BytesIO(img_data))

                # 调整大小
                img.thumbnail(self.max_image_size, Image.Resampling.LANCZOS)

                # 保存缩略图
                thumbnail_data = io.BytesIO()
                img.save(thumbnail_data, format="JPEG", quality=85)
                thumbnail_data.seek(0)

                # 保存文件
                thumbnail_filename = f"thumb_page_{page_num + 1}.jpg"
                thumbnail_url = self._save_thumbnail_file(
                    thumbnail_data.getvalue(), resource, thumbnail_filename
                )

                thumbnail_urls.append(thumbnail_url)

            pdf_document.close()
            return thumbnail_urls

        except Exception as e:
            logger.error(f"PDF缩略图生成失败: {e}")
            raise

    def _generate_pptx_thumbnails(
        self, pptx_data: bytes, resource: MultimediaResource
    ) -> List[str]:
        """生成PPTX缩略图"""
        try:
            prs = Presentation(io.BytesIO(pptx_data))
            thumbnail_urls = []

            # 为前几张幻灯片生成缩略图
            slides_to_thumbnail = min(3, len(prs.slides))

            for slide_num in range(slides_to_thumbnail):
                # PPTX缩略图生成比较复杂，这里简化处理
                # 实际项目中可能需要使用专门的库如python-pptx-image
                prs.slides[slide_num]

                # 创建简单的占位符图像
                img = Image.new("RGB", (800, 600), color=(240, 240, 240))

                # 调整大小
                img.thumbnail(self.max_image_size, Image.Resampling.LANCZOS)

                # 保存缩略图
                thumbnail_data = io.BytesIO()
                img.save(thumbnail_data, format="JPEG", quality=85)
                thumbnail_data.seek(0)

                thumbnail_filename = f"slide_{slide_num + 1}_thumb.jpg"
                thumbnail_url = self._save_thumbnail_file(
                    thumbnail_data.getvalue(), resource, thumbnail_filename
                )

                thumbnail_urls.append(thumbnail_url)

            return thumbnail_urls

        except Exception as e:
            logger.error(f"PPTX缩略图生成失败: {e}")
            return []

    def _generate_cover_thumbnail(
        self,
        document_data: bytes,
        doc_format: DocumentFormat,
        resource: MultimediaResource,
    ) -> Optional[str]:
        """生成封面缩略图"""
        try:
            # 创建代表性的封面图像
            img = Image.new("RGB", (800, 600), color=(245, 245, 245))

            # 添加文档类型标识
            from PIL import ImageDraw, ImageFont

            draw = ImageDraw.Draw(img)

            # 添加文本
            format_text = doc_format.value.upper()
            try:
                # 尝试使用系统字体
                font = ImageFont.truetype("arial.ttf", 48)
            except Exception:
                font = ImageFont.load_default()

            # 计算文本位置
            bbox = draw.textbbox((0, 0), format_text, font=font)
            text_width = bbox[2] - bbox[0]
            text_height = bbox[3] - bbox[1]

            x = (800 - text_width) // 2
            y = (600 - text_height) // 2

            draw.text((x, y), format_text, fill=(100, 100, 100), font=font)

            # 调整大小
            img.thumbnail(self.max_image_size, Image.Resampling.LANCZOS)

            # 保存缩略图
            thumbnail_data = io.BytesIO()
            img.save(thumbnail_data, format="JPEG", quality=85)
            thumbnail_data.seek(0)

            thumbnail_filename = f"cover_thumb.jpg"
            thumbnail_url = self._save_thumbnail_file(
                thumbnail_data.getvalue(), resource, thumbnail_filename
            )

            return thumbnail_url

        except Exception as e:
            logger.error(f"封面缩略图生成失败: {e}")
            return None

    def _save_processed_file(
        self, file_data: bytes, resource: MultimediaResource, file_type: str
    ) -> str:
        """保存处理后的文件"""
        if self.storage_config.storage_type == "s3":
            return self._save_to_s3(file_data, resource, file_type)
        else:
            return self._save_locally(file_data, resource, file_type)

    def _save_thumbnail_file(
        self, file_data: bytes, resource: MultimediaResource, filename: str
    ) -> str:
        """保存缩略图文件"""
        if self.storage_config.storage_type == "s3":
            return self._save_thumbnail_to_s3(file_data, resource, filename)
        else:
            return self._save_thumbnail_locally(file_data, resource, filename)

    def _save_to_s3(
        self, file_data: bytes, resource: MultimediaResource, file_type: str
    ) -> str:
        """保存到S3"""
        try:
            s3_client = self.storage_config.get_s3_client()
            file_key = f"processed/{resource.org_id}/{resource.course_id}/{resource.id}_{file_type}.pdf"

            s3_client.put_object(
                Bucket=self.storage_config.aws_bucket,
                Key=file_key,
                Body=file_data,
                ContentType="application/pdf",
            )

            return f"https://{self.storage_config.aws_bucket}.s3.{self.storage_config.aws_region}.amazonaws.com/{file_key}"

        except Exception as e:
            logger.error(f"S3文件保存失败: {e}")
            raise

    def _save_locally(
        self, file_data: bytes, resource: MultimediaResource, file_type: str
    ) -> str:
        """保存到本地"""
        try:
            processed_dir = os.path.join(
                self.storage_config.local_storage_path,
                "processed",
                str(resource.org_id),
                str(resource.course_id),
            )
            os.makedirs(processed_dir, exist_ok=True)

            filename = f"{resource.id}_{file_type}.pdf"
            file_path = os.path.join(processed_dir, filename)

            with open(file_path, "wb") as f:
                f.write(file_data)

            if self.storage_config.cdn_domain:
                return f"https://{self.storage_config.cdn_domain}/processed/{resource.org_id}/{resource.course_id}/{filename}"
            else:
                return f"/processed/{resource.org_id}/{resource.course_id}/{filename}"

        except Exception as e:
            logger.error(f"本地文件保存失败: {e}")
            raise

    def _save_thumbnail_to_s3(
        self, file_data: bytes, resource: MultimediaResource, filename: str
    ) -> str:
        """保存缩略图到S3"""
        try:
            s3_client = self.storage_config.get_s3_client()
            file_key = f"thumbnails/{resource.org_id}/{resource.course_id}/{resource.id}_{filename}"

            s3_client.put_object(
                Bucket=self.storage_config.aws_bucket,
                Key=file_key,
                Body=file_data,
                ContentType="image/jpeg",
            )

            return f"https://{self.storage_config.aws_bucket}.s3.{self.storage_config.aws_region}.amazonaws.com/{file_key}"

        except Exception as e:
            logger.error(f"S3缩略图保存失败: {e}")
            raise

    def _save_thumbnail_locally(
        self, file_data: bytes, resource: MultimediaResource, filename: str
    ) -> str:
        """保存缩略图到本地"""
        try:
            thumb_dir = os.path.join(
                self.storage_config.local_storage_path,
                "thumbnails",
                str(resource.org_id),
                str(resource.course_id),
            )
            os.makedirs(thumb_dir, exist_ok=True)

            file_path = os.path.join(thumb_dir, f"{resource.id}_{filename}")

            with open(file_path, "wb") as f:
                f.write(file_data)

            if self.storage_config.cdn_domain:
                return f"https://{self.storage_config.cdn_domain}/thumbnails/{resource.org_id}/{resource.course_id}/{resource.id}_{filename}"
            else:
                return f"/thumbnails/{resource.org_id}/{resource.course_id}/{resource.id}_{filename}"

        except Exception as e:
            logger.error(f"本地缩略图保存失败: {e}")
            raise

    def get_supported_formats(self) -> List[DocumentFormat]:
        """获取支持的文档格式"""
        return [
            DocumentFormat.PDF,
            DocumentFormat.MARKDOWN,
            DocumentFormat.PPTX,
            DocumentFormat.DOCX,
            DocumentFormat.TXT,
            DocumentFormat.HTML,
        ]

    def validate_document_file(self, file_data: bytes, file_extension: str) -> bool:
        """验证文档文件格式"""
        try:
            if file_extension.lower() == "pdf":
                return self._validate_pdf_file(file_data)
            elif file_extension.lower() == "docx":
                return self._validate_docx_file(file_data)
            elif file_extension.lower() == "pptx":
                return self._validate_pptx_file(file_data)
            elif file_extension.lower() in ["md", "markdown"]:
                return self._validate_markdown_file(file_data)
            elif file_extension.lower() == "txt":
                return self._validate_txt_file(file_data)
            elif file_extension.lower() == "html":
                return self._validate_html_file(file_data)
            else:
                return False

        except Exception as e:
            logger.error(f"文档文件验证失败: {e}")
            return False

    def _validate_pdf_file(self, file_data: bytes) -> bool:
        """验证PDF文件"""
        return file_data.startswith(b"%PDF")

    def _validate_docx_file(self, file_data: bytes) -> bool:
        """验证DOCX文件"""
        return file_data.startswith(b"PK") and b"[Content_Types].xml" in file_data

    def _validate_pptx_file(self, file_data: bytes) -> bool:
        """验证PPTX文件"""
        if not file_data.startswith(b"PK"):
            return False
        try:
            import zipfile

            with zipfile.ZipFile(io.BytesIO(file_data)) as zf:
                return "[Content_Types].xml" in zf.namelist() and any(
                    name.startswith("ppt/") for name in zf.namelist()
                )
        except Exception:
            return False

    def _validate_markdown_file(self, file_data: bytes) -> bool:
        """验证Markdown文件"""
        try:
            text = file_data.decode("utf-8", errors="ignore")
            # 简单检查是否包含常见的Markdown元素
            markdown_indicators = ["#", "*", "_", "`", "[", "]", "(", ")"]
            return any(indicator in text for indicator in markdown_indicators)
        except Exception:
            return False

    def _validate_txt_file(self, file_data: bytes) -> bool:
        """验证TXT文件"""
        try:
            file_data.decode("utf-8")
            return True
        except Exception:
            return False

    def _validate_html_file(self, file_data: bytes) -> bool:
        """验证HTML文件"""
        try:
            text = file_data.decode("utf-8", errors="ignore").lower()
            return "<html" in text or "<!doctype" in text
        except Exception:
            return False
